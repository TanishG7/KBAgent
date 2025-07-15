# preprocess.py

import os
import io
from pathlib import Path
from typing import Optional
import fitz  #pymupdf
from pptx import Presentation
import docx
from PIL import Image
import pytesseract
import logging
import re

from config import Config
from logger import JSONLogger

# Set tesseract path from config
pytesseract.pytesseract.tesseract_cmd = Config.TESSERACT_CMD

class DocumentProcessor:
    """Handles document processing with comprehensive logging"""
    
    def __init__(self, logger: JSONLogger):
        self.logger = logger
    
    def _clean_text(self, text: str) -> str:
        """Clean text by removing excessive newlines and whitespace"""
        # Normalize line breaks
        text = text.replace('\r\n', '\n').replace('\r', ' ')
        # Replace multiple newlines with single space
        text = re.sub(r'\n+', ' ', text)
        # Remove extra whitespace
        text = re.sub(r'[ \t]+', ' ', text).strip()
        return text
    
    def process_file(self, filepath: str, output_dir: str = None) -> Optional[str]:
        """
        Process a file and convert it to text format
        
        Args:
            filepath: Path to the input file
            output_dir: Output directory for processed text files
            
        Returns:
            Path to the output text file if successful, None otherwise
        """
        MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    
        file_size = os.path.getsize(filepath)
        if file_size > MAX_FILE_SIZE:
            error_msg = f"File too large ({file_size} bytes)"
            self.logger.log_file_processing_error(filename, str(filepath), file_ext, error_msg)
            return None
    
        if output_dir is None:
            output_dir = Config.DOCS_FOLDER
            
        # Extract file information
        file_path = Path(filepath)
        filename = file_path.name
        file_stem = file_path.stem
        file_ext = file_path.suffix.lower()
        output_txt = os.path.join(output_dir, f"{file_stem}.txt")
        
        try:
            # Process based on file type
            if file_ext == ".pdf":
                text = self._extract_pdf_text(filepath)
            elif file_ext in [".pptx", ".ppt"]:
                if file_ext == ".ppt":
                    try:
                        from subprocess import run
                        converted_path = filepath.replace(".ppt", ".pptx")
                        run(["libreoffice", "--headless", "--convert-to", "pptx", filepath, "--outdir", os.path.dirname(filepath)], check=True)
                        filepath = converted_path  # Use converted pptx
                        file_ext = ".pptx"
                        filename = Path(filepath).name
                    except Exception as conv_err:
                        error_msg = f"Conversion failed for .ppt: {str(conv_err)}"
                        self.logger.log_file_processing_error(filename, str(filepath), file_ext, error_msg)
                        return None
                text = self._extract_pptx_text(filepath)
            elif file_ext == ".docx":
                text = self._extract_docx_text(filepath)
            elif file_ext == ".txt":
                with open(filepath, 'r', encoding='utf-8') as f:
                    text = f.read()
            elif file_ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
                text = self._extract_image_text(filepath)
            else:
                error_msg = f"Unsupported file type: {file_ext}"
                self.logger.log_file_processing_error(filename, str(filepath), file_ext, error_msg)
                return None
            
            text = self._clean_text(text)
            
            # Calculate text statistics
            text_length = len(text)
            word_count = len(text.split()) if text else 0
            
            # Write processed text to file
            with open(output_txt, 'w', encoding='utf-8') as f:
                f.write(text)
            
            return output_txt
            
        except Exception as e:
            self.logger.log_file_processing_error(filename, str(filepath), file_ext, str(e))
            return None
    
    def _extract_pdf_text(self, path: str) -> str:
        """Extract text from PDF file using pymupdf with OCR fallback"""
        text = []
        
        try:
            doc = fitz.open(path)
            
            for page in doc:
                # Extract text directly from PDF
                page_text = page.get_text()
                if page_text.strip():
                    text.append(page_text)
                
                # Extract text from images using OCR
                for img in page.get_images():
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        img_bytes = pix.tobytes("png")
                        img_pil = Image.open(io.BytesIO(img_bytes))
                        ocr_text = pytesseract.image_to_string(img_pil)
                        if ocr_text.strip():
                            text.append(ocr_text)
                    except Exception as img_e:
                        logging.warning(f"Image OCR failed on page {page.number}: {str(img_e)}")
                        continue
            
            return "\n".join(text)
            
        except Exception as e:
            logging.error(f"PDF processing failed: {str(e)}")
            raise
    
    def _extract_pptx_text(self, path: str) -> str:
        """Extract text from PowerPoint file"""
        prs = Presentation(path)
        text = []
        
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
                
                # Extract text from images using OCR
                try:
                    if shape.shape_type == 13:  # Picture shape type
                        img_bytes = shape.image.blob
                        with Image.open(io.BytesIO(img_bytes)) as img:
                            img_text = pytesseract.image_to_string(img)
                            if img_text.strip():
                                text.append(img_text)
                except Exception as e:
                    logging.warning(f"Slide {slide_num + 1} image OCR failed: {str(e)}")
        
        return "\n".join(text)
    
    def _extract_docx_text(self, path: str) -> str:
        """Extract text from Word document"""
        doc = docx.Document(path)
        return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
    
    def _extract_image_text(self, path: str) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(path)
            return pytesseract.image_to_string(image)
        except Exception as e:
            logging.error(f"Image OCR failed: {str(e)}")
            return ""

# Legacy function for backward compatibility
def process_file(filepath, output_dir="docs"):
    """Legacy function - creates a temporary logger for single file processing"""
    from datetime import datetime
    log_path = f"temp_process_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
    temp_logger = JSONLogger(log_path, "temp_processor")
    processor = DocumentProcessor(temp_logger)
    return processor.process_file(filepath, output_dir)